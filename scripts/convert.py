#!/usr/bin/env python3
"""
pptx/ 폴더의 파일들을 PNG 슬라이드 이미지로 변환합니다.

지원 형식:
  - .pptx / .ppt  → LibreOffice로 PDF 변환 후 PNG 렌더링
                    (+ python-pptx 로 링크/동영상 추출, .pptx만 해당)
  - .pdf          → 바로 PNG 렌더링 (LibreOffice 변환 단계 생략)
                    (+ PyMuPDF 로 PDF 안의 링크 추출)

추출된 "링크(같은 발표자료 내 슬라이드 이동 포함)"와 "삽입된 동영상"은
이미지 위에 겹쳐 올릴 수 있는 좌표(overlay) 정보로 manifest.json 에 기록됩니다.

결과:
  docs/slides/<deck-id>/slide_0001.png ...
  docs/slides/<deck-id>/media/*.mp4        (pptx에서 추출된 동영상)
  docs/manifest.json

이 스크립트는 사람이 직접 실행하지 않고, .github/workflows/convert.yml 이
pptx/ 폴더에 변경이 생길 때마다 자동으로 실행합니다.
"""
import json
import re
import shutil
import subprocess
import sys
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.action import PP_ACTION
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC_DIR = ROOT / 'pptx'
DOCS_DIR = ROOT / 'docs'
SLIDES_DIR = DOCS_DIR / 'slides'
DPI = 288  # 3x 고화질 (큰 화면/4K 모니터 대비)

NAME_SUFFIX_RE = re.compile(r'\.(pptx|ppt|pdf)$', re.I)


def slugify(name: str) -> str:
    s = NAME_SUFFIX_RE.sub('', name)
    s = re.sub(r'[^a-zA-Z0-9가-힣_-]+', '-', s).strip('-')
    return s or 'deck'


def title_of(name: str) -> str:
    return NAME_SUFFIX_RE.sub('', name)


# ── LibreOffice: pptx/ppt → pdf ─────────────────────────────────
def convert_to_pdf(src: Path, out_dir: Path) -> Path:
    r = subprocess.run(
        ['libreoffice', '--headless', '--norestore',
         '--convert-to', 'pdf', '--outdir', str(out_dir), str(src)],
        capture_output=True, text=True, timeout=300
    )
    if r.returncode != 0:
        raise RuntimeError(f'{src.name} 변환 실패:\n{r.stderr[-500:]}')
    pdf = out_dir / (src.stem + '.pdf')
    if not pdf.exists():
        pdfs = list(out_dir.glob('*.pdf'))
        if not pdfs:
            raise RuntimeError(f'{src.name}: PDF가 생성되지 않았습니다')
        pdf = pdfs[-1]
    return pdf


# ── PDF → PNG ────────────────────────────────────────────────────
def render_pngs(pdf_path: Path, out_dir: Path, deck_id: str):
    doc = fitz.open(str(pdf_path))
    mat = fitz.Matrix(DPI / 72, DPI / 72)
    r0 = doc[0].rect
    ratio = round(r0.width / r0.height, 6)
    urls = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB, alpha=False)
        fname = f'slide_{i + 1:04d}.png'
        pix.save(str(out_dir / fname))
        urls.append(f'slides/{deck_id}/{fname}')
    n = doc.page_count
    doc.close()
    return urls, ratio, n


# ── PDF 자체의 링크 추출 (PDF로 직접 올린 경우) ─────────────────
def extract_pdf_overlays(pdf_path: Path):
    doc = fitz.open(str(pdf_path))
    result = []
    for page in doc:
        rect = page.rect
        overlays = []
        for link in page.get_links():
            r = link.get('from')
            if not r or not rect.width or not rect.height:
                continue
            box = {
                'x': round(r.x0 / rect.width * 100, 3),
                'y': round(r.y0 / rect.height * 100, 3),
                'w': round((r.x1 - r.x0) / rect.width * 100, 3),
                'h': round((r.y1 - r.y0) / rect.height * 100, 3),
            }
            kind = link.get('kind')
            if kind == fitz.LINK_GOTO:
                tp = link.get('page')
                if tp is not None and tp >= 0:
                    overlays.append({'type': 'link', 'target': {'kind': 'internal', 'slide': tp + 1}, **box})
            elif kind == fitz.LINK_URI:
                uri = link.get('uri')
                if uri:
                    overlays.append({'type': 'link', 'target': {'kind': 'external', 'url': uri}, **box})
        result.append(overlays)
    doc.close()
    return result


# ── PPTX 자체의 링크 / 동영상 추출 ───────────────────────────────
def _slide_id_map(prs):
    return {slide.slide_id: idx for idx, slide in enumerate(prs.slides)}


def _resolve_click_action(shape, slide_id_map, n_slides):
    """도형 전체에 걸린 '실행 단추/작업' 링크 (Insert > Action)"""
    try:
        ca = shape.click_action
    except Exception:
        return None
    try:
        if ca.action == PP_ACTION.HYPERLINK:
            if ca.hyperlink and ca.hyperlink.address:
                return {'kind': 'external', 'url': ca.hyperlink.address}
            if ca.target_slide is not None:
                idx = slide_id_map.get(ca.target_slide.slide_id)
                if idx is not None:
                    return {'kind': 'internal', 'slide': idx + 1}
        elif ca.action == PP_ACTION.FIRST_SLIDE:
            return {'kind': 'internal', 'slide': 1}
        elif ca.action == PP_ACTION.LAST_SLIDE:
            return {'kind': 'internal', 'slide': n_slides}
        elif ca.action == PP_ACTION.NEXT_SLIDE:
            return {'kind': 'relative', 'delta': 1}
        elif ca.action == PP_ACTION.PREVIOUS_SLIDE:
            return {'kind': 'relative', 'delta': -1}
    except Exception:
        pass
    return None


def _resolve_run_hyperlinks(shape, slide_part, slide_part_map):
    """텍스트 안의 특정 단어/문구에 걸린 하이퍼링크 (있으면 도형 전체를 클릭 영역으로 사용)"""
    if not shape.has_text_frame:
        return None
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                hlink = run.hyperlink
                if hlink is None:
                    continue
                if hlink.address:
                    return {'kind': 'external', 'url': hlink.address}
                raw = getattr(hlink, '_hlink', None)
                if raw is None:
                    continue
                rid = raw.get(qn('r:id'))
                if not rid:
                    continue
                try:
                    target_part = slide_part.related_part(rid)
                    idx = slide_part_map.get(id(target_part))
                    if idx is not None:
                        return {'kind': 'internal', 'slide': idx + 1}
                except Exception:
                    pass
    except Exception:
        pass
    return None


MAX_VIDEO_MB = 90  # GitHub 100MB 제한보다 여유있게 잡은 상한선


def _compress_video(raw_bytes: bytes, out_path: Path) -> bool:
    """추출한 동영상을 웹 재생에 적합한 크기로 압축(mp4/h264). 성공하면 True."""
    tmp_in = out_path.with_name(out_path.stem + '_src.tmp')
    tmp_in.write_bytes(raw_bytes)
    try:
        r = subprocess.run([
            'ffmpeg', '-y', '-i', str(tmp_in),
            '-vf', "scale='min(1280,iw)':-2",
            '-c:v', 'libx264', '-preset', 'veryfast', '-crf', '28',
            '-c:a', 'aac', '-b:a', '128k',
            '-movflags', '+faststart',
            str(out_path)
        ], capture_output=True, text=True, timeout=600)
        return r.returncode == 0 and out_path.exists()
    except Exception:
        return False
    finally:
        tmp_in.unlink(missing_ok=True)


def _find_video_rid(xml: str):
    """도형 XML에서 동영상 관계 rId를 찾는다.
    PowerPoint는 실제로 삽입(embed)된 동영상도 태그 이름을 'r:link'로 쓰는 경우가
    많아서(마이크로소프트 특유의 헷갈리는 네이밍), 태그 이름만으로는 내장 여부를
    판단할 수 없다. 여기서는 rId 후보만 찾고, 실제 내장 여부는 rels 에서 확인한다.
    p14:media(r:embed) 가 있으면 그게 항상 우선 (PowerPoint 2010+ 호환용 사본).
    """
    m = re.search(r'<[a-zA-Z0-9]*:media\b[^>]*r:embed="([^"]+)"', xml)
    if m:
        return m.group(1)
    m = re.search(r'videoFile[^>]*r:(?:embed|link)="([^"]+)"', xml)
    if m:
        return m.group(1)
    return None


def _extract_video(shape, slide_part, media_dir, media_prefix, counter):
    """도형 XML에서 삽입된 동영상을 찾아 압축 후 저장. 없거나(=외부 연결) 실패하면 None."""
    try:
        xml = shape._element.xml
    except Exception:
        return None
    rid = _find_video_rid(xml)
    if not rid:
        return None
    try:
        rel = slide_part.rels[rid]
        if rel.is_external:
            return None  # 실제로 외부 URL과 연결된 동영상 (내장 아님) — 추출 불가
        part = rel.target_part
        blob = part.blob
        media_dir.mkdir(exist_ok=True)
        fname = f'{media_prefix}_{counter}.mp4'
        out_path = media_dir / fname
        if not _compress_video(blob, out_path):
            out_path.unlink(missing_ok=True)
            print(f'  [알림] 동영상 압축 실패 (건너뜀)')
            return None
        size_mb = out_path.stat().st_size / (1024 * 1024)
        if size_mb > MAX_VIDEO_MB:
            print(f'  [알림] 동영상이 압축 후에도 {size_mb:.0f}MB라 제외됨 (깃허브 용량 제한)')
            out_path.unlink(missing_ok=True)
            return None
        return f'media/{fname}'
    except Exception as e:
        print(f'  [알림] 동영상 추출 실패: {e}')
        return None


def extract_pptx_overlays(pptx_path: Path, out_dir: Path, deck_id: str, n_slides: int):
    """슬라이드별 [{type, x,y,w,h(%), ...}] 리스트 반환. 실패해도 빈 리스트로 조용히 넘어감."""
    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        print(f'  [알림] 링크/동영상 추출 건너뜀 (python-pptx 열기 실패: {e})')
        return [[] for _ in range(n_slides)]

    sw, sh = prs.slide_width, prs.slide_height
    if not sw or not sh:
        return [[] for _ in range(n_slides)]

    slide_id_map = _slide_id_map(prs)
    slide_part_map = {id(slide.part): idx for idx, slide in enumerate(prs.slides)}
    media_dir = out_dir / 'media'
    result = []
    vcount = 0

    for si, slide in enumerate(prs.slides):
        overlays = []
        for shape in slide.shapes:
            if shape.left is None or shape.top is None or not shape.width or not shape.height:
                continue
            box = {
                'x': round(shape.left / sw * 100, 3),
                'y': round(shape.top / sh * 100, 3),
                'w': round(shape.width / sw * 100, 3),
                'h': round(shape.height / sh * 100, 3),
            }

            vcount += 1
            video_url = _extract_video(shape, slide.part, media_dir, f's{si + 1}', vcount)
            if video_url:
                overlays.append({'type': 'video', 'src': video_url, **box})
                continue  # 동영상 도형이면 링크 체크는 생략

            target = _resolve_click_action(shape, slide_id_map, n_slides)
            if target is None:
                target = _resolve_run_hyperlinks(shape, slide.part, slide_part_map)
            if target:
                overlays.append({'type': 'link', 'target': target, **box})

        result.append(overlays)

    if media_dir.exists() and not any(media_dir.iterdir()):
        media_dir.rmdir()

    return result


# ── 메인 ─────────────────────────────────────────────────────────
def is_lfs_pointer(path: Path) -> bool:
    """Git LFS가 실제 파일 대신 '포인터(껍데기)' 상태로 남아있는지 확인.
    (100~200바이트 정도의 작은 텍스트 파일로, 실제 파일 내용이 아직 안 받아진 상태)"""
    try:
        if path.stat().st_size > 1024:  # 정상 파일은 훨씬 큼
            return False
        with open(path, 'rb') as f:
            head = f.read(60)
        return b'git-lfs.github.com' in head
    except Exception:
        return False


def collect_sources():
    items = []
    for src in SRC_DIR.glob('*'):
        if src.is_file() and src.suffix.lower() in ('.pptx', '.ppt', '.pdf'):
            items.append(src)
    return sorted(items, key=lambda p: p.name.lower())


def main():
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    sources = collect_sources()

    decks = []
    skipped_ids = set()
    for src in sources:
        ext = src.suffix.lower()
        deck_id = slugify(src.name)

        if is_lfs_pointer(src):
            print(f'[오류] {src.name}: Git LFS 파일이 아직 실제 내용으로 받아지지 않았습니다 '
                  f'(포인터 상태). 워크플로우의 "git lfs pull" 단계를 확인하세요.', file=sys.stderr)
            skipped_ids.add(deck_id)
            continue

        print(f'[변환] {src.name} -> {deck_id}')
        tmp_dir = ROOT / f'_tmp_{deck_id}'
        tmp_dir.mkdir(exist_ok=True)
        try:
            out_dir = SLIDES_DIR / deck_id
            if out_dir.exists():
                shutil.rmtree(out_dir)
            out_dir.mkdir(parents=True)

            if ext == '.pdf':
                pdf = src  # 이미 PDF라 변환 단계 생략
            else:
                pdf = convert_to_pdf(src, tmp_dir)

            urls, ratio, n = render_pngs(pdf, out_dir, deck_id)

            if ext == '.pptx':  # python-pptx 는 구형 .ppt 미지원
                overlays_per_slide = extract_pptx_overlays(src, out_dir, deck_id, n)
            elif ext == '.pdf':
                overlays_per_slide = extract_pdf_overlays(pdf)
            else:
                overlays_per_slide = [[] for _ in range(n)]

            if len(overlays_per_slide) != n:  # 페이지 수 불일치 시 안전하게 무시
                overlays_per_slide = [[] for _ in range(n)]

            slides = [
                {'index': i + 1, 'url': urls[i], 'overlays': overlays_per_slide[i]}
                for i in range(n)
            ]
            n_links = sum(1 for s in slides for o in s['overlays'] if o['type'] == 'link')
            n_videos = sum(1 for s in slides for o in s['overlays'] if o['type'] == 'video')

            decks.append({
                'id': deck_id,
                'filename': src.name,
                'title': title_of(src.name),
                'source_type': ext.lstrip('.'),
                'slide_count': n,
                'aspect_ratio': ratio,
                'slides': slides,
            })
            print(f'  -> {n}장 완료 (링크 {n_links}개, 동영상 {n_videos}개)')
        except Exception as e:
            print(f'[오류] {src.name}: {e}', file=sys.stderr)
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    # pptx/ 에서 삭제된 파일의 이미지 폴더는 정리
    # (단, LFS 못 받아온 파일은 '삭제된 것'이 아니라 '일시적으로 실패한 것'이므로
    #  기존에 만들어둔 이미지가 있다면 지우지 않고 그대로 둡니다)
    valid_ids = {d['id'] for d in decks} | skipped_ids
    if SLIDES_DIR.exists():
        for child in SLIDES_DIR.iterdir():
            if child.is_dir() and child.name not in valid_ids:
                shutil.rmtree(child, ignore_errors=True)

    manifest = {'decks': decks}
    (DOCS_DIR / 'manifest.json').write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2)
    )
    print(f'\n완료: 총 {len(decks)}개 발표자료 변환됨')


if __name__ == '__main__':
    main()
